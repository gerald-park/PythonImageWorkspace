from pptx import Presentation
import pandas as pd

MAX_ITEM_NO = 12


def getTextInRowCol(table, row, col) : 
    text = ''
    for paragraph in table.cell(row,col).text_frame.paragraphs:
        for run in paragraph.runs:
            text += run.text 
    return text

def GetItemInTable(table, idx):        
    OP_ROW = 1 
    OP_COL = 3
    FIRST_ROW = 3
    LEFT_ITEM_COL = 3 
    LEFT_DES_COL = 9
    LEFT_QTY_COL = 13
    RIGHT_ITEM_COL = 15 
    RIGHT_DES_COL = 18
    RIGHT_QTY_COL = 21
    
    right_first_item_no = int(MAX_ITEM_NO / 2)
    
    if idx < right_first_item_no :
        op_num = getTextInRowCol(table, OP_ROW, OP_COL)
        item_no = getTextInRowCol(table, FIRST_ROW + idx, LEFT_ITEM_COL)
        des = getTextInRowCol(table, FIRST_ROW + idx, LEFT_DES_COL )
        qty = getTextInRowCol(table, FIRST_ROW + idx, LEFT_QTY_COL)

    else:
        op_num = getTextInRowCol(table, OP_ROW, OP_COL)
        item_no = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_ITEM_COL)
        des = getTextInRowCol(table, FIRST_ROW  + idx - right_first_item_no, RIGHT_DES_COL)
        qty = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_QTY_COL)
    
    item_list = [op_num, item_no, des, qty]         
    return item_list

def GetItemsInSlide(slide):
    items = []
    if slide.shapes[0].has_table:
        table = slide.shapes[0].table
        for idx in range(0, MAX_ITEM_NO):
            item = GetItemInTable(table, idx)
            if item[3].isnumeric() and item[2]: # Description에 데이터가 있고 qty가 숫자이면 데이터 인정 
                items.append(item)
        return items

df = pd.DataFrame(columns=['OperationStep','Man.Item.No','Description','Qty'])
prs = Presentation("../sop.pptx")
for slide in prs.slides:
    items = GetItemsInSlide(slide) 
    if items : # items list 에 데이터가 있으면 추가 
        for item in items:
            df.loc[len(df)] = item
            
f_name = 'sop_pptx.csv'
df.to_csv(f_name,encoding='utf-8-sig', index=False, mode='w', header=True)