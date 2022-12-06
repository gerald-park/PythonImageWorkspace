from pptx import Presentation
import pandas as pd
import tkinter.ttk as ttk
import tkinter.messagebox as msgbox
from tkinter import * # __all__
from tkinter import filedialog

class Point: 
    def __init__(self, row, col):
        self.row = row
        self.col = col
        
    row: int = None
    col: int = None
        
class SOPDataPoints:    
    op_num : Point = None
    item_no : Point = None
    desciption : Point = None
    qty : Point = None   

def getTablesInSlide(prs, sld_no):
    slide_layout = prs.slide_layouts[0]
    shapes = slide_layout.shapes
    slide = prs.slides[sld_no]
    if slide.shapes[0].has_table:
        return slide.shapes[0].table

def getTextInRowCol(tb, row, col) : 
    text = ''
    if not tb.cell(row,col):
        pass
    else :
        for paragraph in tb.cell(row,col).text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text 
    return text

def GetItemsList(tb, i):
    item_no = []
    op_num = Point(1,3)
    item_no = Point(3,3)
    desciption = Point(3,9)
    qty = Point(3,13)
    try:
        if i < 6 : # left
            op_num = getTextInRowCol(tb, op_num.row, op_num.col)
            item_no = getTextInRowCol(tb, item_no.row + i, item_no.col)
            des = getTextInRowCol(tb, desciption.row  + i, desciption.col )
            qty = getTextInRowCol(tb, qty.row + i, qty.col)

        else: # right
            op_num = getTextInRowCol(tb,op_num.row, op_num.col)
            item_no = getTextInRowCol(tb,item_no.row + i - 6, item_no.col + 12)
            des = getTextInRowCol(tb,desciption.row  + i - 6, desciption.col + 9)
            qty = getTextInRowCol(tb,qty.row + i - 6, qty.col + 8)
        item_list = [op_num, item_no, des, qty]   
        return item_list
    except:
        pass    

sopfile = ''

def add_file():
    global sopfile
    global file_name
    try:
        sopfile = filedialog.askopenfilename(initialdir=r'C:\Users\ParkGY\DocumentsCFLTYanadoo\DT Academy\SOPMeterial',title="select a file",
                                            filetypes =(("Perenstaion","*.pptx"),
                                            ("all files","*.*")))
    except:
        pass
    
def convert():        
    prs = Presentation(sopfile)  
    df = pd.DataFrame(columns=['OperationStep','Man.Item.No','Description','Qty'])
    for sld_no in range(0, len(prs.slides)):
        tb = getTablesInSlide(prs, sld_no)
        if not tb:
            continue
        for i in range(0, 12):
            it = GetItemsList(tb, i)
            if it and not it[3].isnumeric():
                continue
            else:
                df.loc[len(df.index)] = it
            
    f_name = sopfile.split('.')[0] + '.csv' 
    df.to_csv(f_name, encoding='utf-8-sig', index=False, mode='w', header=True)   

root = Tk()
root.title("PPTX SOP Converter")
root.geometry("300x150")

# Frame 
file_frame = Frame(root)
file_frame.pack(fill="x", padx=5, pady=5) # 간격 띄우기

btn_add_file = Button(file_frame, padx=5, pady=5, text="Open SOP", width=20, command=add_file)
btn_add_file.pack(side="top")

btn_start = Button(file_frame, padx=5, pady=5, text="Convert", width=20, command=convert)
btn_start.pack(side="top", padx=5, pady=5)
# Runs
root.mainloop()