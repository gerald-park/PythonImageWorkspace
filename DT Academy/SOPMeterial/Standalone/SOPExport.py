from pptx import Presentation
import pandas as pd


class Point: 
    def __init__(self, row, col):
        self.row = row
        self.col = col
        
    row: int = None
    col: int = None
        
class SOPDataPoints:    
    op_num : Point = None
    item_no : Point = None
    desciption : Point = None
    qty : Point = None   

def getTablesInSlide(prs, sld_no):
    slide_layout = prs.slide_layouts[0]
    shapes = slide_layout.shapes
    slide = prs.slides[sld_no]
    if slide.shapes[0].has_table:
        return slide.shapes[0].table

def getTextInRowCol(tb, row, col) : 
    text = ''
    if not tb.cell(row,col):
        pass
    else :
        for paragraph in tb.cell(row,col).text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text 
    return text

def GetItemsList(tbl, cells, i):
    item_no = []
    try:
        if i < 6 : # left
            op_num = getTextInRowCol(tb, cells.op_num.row, cells.op_num.col)
            item_no = getTextInRowCol(tb, cells.item_no.row + i, cells.item_no.col)
            des = getTextInRowCol(tb, cells.desciption.row  + i, cells.desciption.col )
            qty = getTextInRowCol(tb, cells.qty.row + i, cells.qty.col)

        else: # right
            op_num = getTextInRowCol(tb,cells.op_num.row, cells.op_num.col)
            item_no = getTextInRowCol(tb,cells.item_no.row + i - 6, cells.item_no.col + 12)
            des = getTextInRowCol(tb,cells.desciption.row  + i - 6, cells.desciption.col + 9)
            qty = getTextInRowCol(tb,cells.qty.row + i - 6, cells.qty.col + 8)
        item_list = [op_num, item_no, des, qty]   
        return item_list
    except:
        pass     


# define first items
prs = Presentation('sop.pptx')
item = SOPDataPoints()
item.op_num = Point(1,3)
item.item_no = Point(3,3)
item.desciption = Point(3,9)
item.qty = Point(3,13)

df = pd.DataFrame(columns=['OperationStep','Man.Item.No','Description','Qty'])
for sld_no in range(0, len(prs.slides)):
    tb = getTablesInSlide(prs, sld_no)
    if not tb:
        continue
    for i in range(0, 12):
        it = GetItemsList(tb,item, i)
        if it and not it[3].isnumeric():
            continue
        else:
            df.loc[len(df.index)] = it

f_name = 'sop_material.csv'
df.to_csv(f_name,encoding='utf-8-sig', index=False, mode='w', header=True)